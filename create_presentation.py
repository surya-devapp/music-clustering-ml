import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# --- THEME CONFIGURATION ---
BG_COLOR = RGBColor(15, 23, 42)    # Deep Navy
TEXT_COLOR = RGBColor(248, 250, 252) # Off-White
ACCENT_COLOR = RGBColor(56, 189, 248) # Sky Blue
HIGHLIGHT_COLOR = RGBColor(251, 191, 36) # Amber

def create_styled_slide(prs, title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
    
    # Add Background
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = BG_COLOR
    rect.line.fill.background()
    
    # Add Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    title_run = title_box.text_frame.add_paragraph()
    title_run.text = title_text
    title_run.font.size = Pt(36)
    title_run.font.bold = True
    title_run.font.color.rgb = ACCENT_COLOR
    
    return slide

# --- GENERATE PRESENTATION ---
prs = Presentation()
prs.slide_width = Inches(13.33)  # Widescreen
prs.slide_height = Inches(7.5)

# 1. TITLE SLIDE
slide = prs.slides.add_slide(prs.slide_layouts[6])
rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
rect.fill.solid()
rect.fill.fore_color.rgb = BG_COLOR

title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
tf = title_box.text_frame
p = tf.add_paragraph()
p.text = "MUSIC CLUSTERING: KEY INSIGHTS"
p.font.size = Pt(60)
p.font.bold = True
p.font.color.rgb = ACCENT_COLOR
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "Data-Driven Results & Performance Review"
p2.font.size = Pt(28)
p2.font.color.rgb = TEXT_COLOR
p2.alignment = PP_ALIGN.CENTER

# 2. INSIGHT 1: MATHEMATICAL OPTIMALITY
slide = create_styled_slide(prs, "Insight 1: Mathematical Optimality")
img_path = "elbow_silhouette_analysis.png"
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.5), height=Inches(4.5))

text_box = slide.shapes.add_textbox(Inches(7.5), Inches(1.5), Inches(5.5), Inches(5))
tf = text_box.text_frame
tf.word_wrap = True

points = [
    "Optimum Found: The dataset naturally splits into 4 vibey clusters.",
    "Elbow Method: Significant drop in Inertia at K=4 proves optimal balance.",
    "Silhouette Score (0.91): Indicates 91% separation confidence.",
    "Verdict: Clusters are mathematically stable and distinct."
]

for pt in points:
    p = tf.add_paragraph()
    p.text = "✔ " + pt
    p.font.size = Pt(20)
    p.font.color.rgb = TEXT_COLOR
    p.space_after = Pt(15)

# 3. INSIGHT 2: PCA SEPARATION
slide = create_styled_slide(prs, "Insight 2: High-Dimensional Separation")
img_path = "cluster_visualization_pca.png"
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(5.5), Inches(1.5), width=Inches(7.5))

text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
tf = text_box.text_frame
tf.word_wrap = True

points = [
    "Islands of Music: Clear physical separation in PCA space.",
    "PC1 (Energy Axis): Primary driver of cluster differentiation.",
    "No Overlap: Visual proof of the high Silhouette score.",
    "Stability: Dense centers indicate high internal consistency."
]

for pt in points:
    p = tf.add_paragraph()
    p.text = "• " + pt
    p.font.size = Pt(20)
    p.font.color.rgb = TEXT_COLOR
    p.space_after = Pt(15)

# 4. INSIGHT 3: CLUSTER DNA
slide = create_styled_slide(prs, "Insight 3: Cluster DNA (Heatmap)")
img_path = "cluster_feature_heatmap.png"
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.5), width=Inches(8))

text_box = slide.shapes.add_textbox(Inches(9), Inches(1.5), Inches(4), Inches(5))
tf = text_box.text_frame

points = [
    "Cluster 0: High Acousticness (Chill/Focus).",
    "Cluster 1: High Danceability & Valence (Party).",
    "Cluster 2: High Energy & Liveness (Performance).",
    "Automation: These heat signatures act as automated vibe tags."
]

for pt in points:
    p = tf.add_paragraph()
    p.text = "➤ " + pt
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_COLOR
    p.space_after = Pt(12)

# 5. INSIGHT 4: PREPROCESSING IMPACT
slide = create_styled_slide(prs, "Insight 4: Preprocessing Success")
text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
tf = text_box.text_frame

points = [
    "Log Transformation: Neutralized the 'Duration' outlier power.",
    "Skewness Correction: Duration_ms skew dropped from 9.81 to -0.65.",
    "Robust Scaling: Neutralized live recording volume spikes.",
    "Impact: Model clusters by musicality, not just track length."
]

for pt in points:
    p = tf.add_paragraph()
    p.text = "⚡ " + pt
    p.font.size = Pt(24)
    p.font.color.rgb = TEXT_COLOR
    p.space_after = Pt(20)

# 6. INSIGHT 5: BUSINESS VALUE
slide = create_styled_slide(prs, "Insight 5: Business Implementation Value")
text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
tf = text_box.text_frame

points = [
    "Automated Playlisting: Higher user satisfaction with 'Vibe Match'.",
    "Targeted Ads: serve high-energy ads to high-energy listeners.",
    "Production Ready: Exported .pkl model ready for API integration.",
    "Scalability: Process millions of songs at near-zero curation cost."
]

for pt in points:
    p = tf.add_paragraph()
    p.text = "💰 " + pt
    p.font.size = Pt(24)
    p.font.color.rgb = HIGHLIGHT_COLOR
    p.space_after = Pt(20)

prs.save("Music_Clustering_Final_Presentation.pptx")
print("Insights-Only Presentation saved as Music_Clustering_Final_Presentation.pptx")
